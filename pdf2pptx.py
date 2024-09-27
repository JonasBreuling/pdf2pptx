import os, sys
import argparse

from pdf2image import convert_from_path
from pptx import Presentation
from io import BytesIO

# parse input arguments
parser = argparse.ArgumentParser(description="Convert pdf file (e.g. LaTeX beamer or CorelDraw presentation) to pptx presentation.")
parser.add_argument('pdf', metavar='pdf', type=str, help='relative path to pdf file')
parser.add_argument('-out', metavar='out', type=str, help="relative path to output presentation (ignores extensions and appends .pptx)")
parser.add_argument('-dpi', metavar='dpi', default=300, help='desired dpi, default is 300')
args = parser.parse_args()

# the pdf we want to convert
pdf_file = args.pdf
print()
print("Converting file: " + pdf_file)
print()

# initialize presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

# use output file if given, otherwise use pdf filename
if args.out is not None:
	out_name = os.path.splitext(args.out)[0]
else:
	out_name = pdf_file.split(".pdf")[0]

# Convert PDF to list of images
print("Starting conversion...")
slideimgs = convert_from_path(pdf_file, args.dpi, fmt='ppm', thread_count=2)
print("...complete.")
print()

# loop over slides
for i, slideimg in enumerate(slideimgs):
	if i % 10 == 0:
		print("Saving slide: " + str(i))

	imagefile = BytesIO()
	slideimg.save(imagefile, format='tiff')
	imagedata = imagefile.getvalue()
	imagefile.seek(0)
	width, height = slideimg.size

	# set slide dimensions
	prs.slide_height = height * 9525
	prs.slide_width = width * 9525

	# add slide
	slide = prs.slides.add_slide(blank_slide_layout)
	pic = slide.shapes.add_picture(imagefile, 0, 0, width=width * 9525, height=height * 9525)

# save Powerpoint presentation
print()
print("Saving file: " + out_name + ".pptx")
prs.save(out_name + '.pptx')
print("Conversion complete. :)")